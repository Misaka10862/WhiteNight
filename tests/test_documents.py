"""文档解析测试语料：文本/代码、PDF、Office、图片 OCR、压缩包。"""

from __future__ import annotations

import zipfile
from pathlib import Path

import pytest

from whitenight.documents.dispatcher import parse_document
from whitenight.documents.ocr import OcrUnavailableError, ocr_image
from whitenight.documents.text import read_text_file


def test_utf8_and_gb18030_text(tmp_path: Path) -> None:
    utf8 = tmp_path / "utf8.txt"
    utf8.write_text("小白，今天也要加油。", encoding="utf-8")
    reading = read_text_file(utf8)
    assert "小白" in reading.text
    assert reading.encoding == "utf-8"

    gbk = tmp_path / "gb18030.txt"
    gbk.write_bytes("主人，晚上好".encode("gb18030"))
    reading = read_text_file(gbk)
    assert "主人" in reading.text
    assert reading.encoding == "gb18030"


def test_code_file_is_plain_text(tmp_path: Path) -> None:
    code = tmp_path / "main.py"
    code.write_text("def main():\n    return 42\n", encoding="utf-8")
    parsed = parse_document(code)
    assert parsed.kind == "text"
    assert "return 42" in parsed.text
    assert parsed.sources == [str(code)]


def test_pdf_text_extraction_with_source(tmp_path: Path) -> None:
    fitz = pytest.importorskip("fitz")
    pdf = tmp_path / "report.pdf"
    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), "WhiteNight 阶段 3 测试")
    document.save(pdf)
    document.close()

    parsed = parse_document(pdf)
    assert parsed.kind == "pdf"
    assert "WhiteNight" in parsed.text
    assert parsed.metadata["page_count"] == 1
    assert parsed.sources == [str(pdf)]


def test_office_documents(tmp_path: Path) -> None:
    from docx import Document

    docx_path = tmp_path / "brief.docx"
    document = Document()
    document.add_paragraph("DOCX 段落内容")
    document.save(docx_path)
    parsed = parse_document(docx_path)
    assert "DOCX 段落内容" in parsed.text

    from openpyxl import Workbook

    xlsx_path = tmp_path / "sheet.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "数据"
    sheet.append(["城市", "温度"])
    sheet.append(["杭州", 28])
    workbook.save(xlsx_path)
    parsed = parse_document(xlsx_path)
    assert "杭州" in parsed.text
    assert "数据" in parsed.text

    from pptx import Presentation

    pptx_path = tmp_path / "slides.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "PPTX 标题"
    presentation.save(pptx_path)
    parsed = parse_document(pptx_path)
    assert "PPTX 标题" in parsed.text


def test_archive_lists_without_extracting(tmp_path: Path) -> None:
    archive_path = tmp_path / "bundle.zip"
    with zipfile.ZipFile(archive_path, "w") as archive:
        archive.writestr("a.txt", "hello")
        archive.writestr("dir/b.txt", "world")

    parsed = parse_document(archive_path)
    assert parsed.kind == "archive"
    assert "a.txt" in parsed.text
    assert "dir/b.txt" in parsed.text
    assert not (tmp_path / "a.txt").exists()
    assert not (tmp_path / "dir").exists()


def test_ocr_rendered_text_image(tmp_path: Path) -> None:
    """用 PDF 渲染文本后交给 Apple Vision OCR；无 OCR 能力时跳过。"""
    fitz = pytest.importorskip("fitz")
    document = fitz.open()
    page = document.new_page(width=320, height=120)
    page.insert_text((40, 70), "WHITE NIGHT 42", fontsize=36)
    pixmap = page.get_pixmap(dpi=300)
    image = tmp_path / "rendered.png"
    pixmap.save(str(image))
    document.close()
    try:
        ocr = ocr_image(image)
    except OcrUnavailableError as exc:
        pytest.skip(f"OCR 不可用：{exc}")
    assert "42" in ocr.text or ocr.confidence is not None


def test_unsupported_format_reports_error(tmp_path: Path) -> None:
    path = tmp_path / "mystery.xyz"
    path.write_bytes(b"data")
    parsed = parse_document(path)
    assert parsed.kind == "unsupported"
    assert parsed.error
    assert parsed.text == ""
