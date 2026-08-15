"""Office 文档解析：不执行宏或嵌入脚本，只提取可见文本。"""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path

_MAX_ROWS = 500
_MAX_COLS = 50


class OfficeParseError(ValueError):
    """Office 文档解析失败。"""


@dataclass
class OfficeReading:
    text: str
    parts: list[str] = field(default_factory=list)
    metadata: dict[str, object] = field(default_factory=dict)


def parse_docx(path: Path, max_chars: int = 200_000) -> OfficeReading:
    try:
        from docx import Document
    except ImportError as exc:
        raise OfficeParseError("缺少 python-docx 依赖") from exc
    try:
        document = Document(str(path))
        parts: list[str] = [paragraph.text for paragraph in document.paragraphs if paragraph.text]
        for table in document.tables:
            for row in table.rows:
                parts.append(" | ".join(cell.text for cell in row.cells))
        text = "\n".join(parts)
        return OfficeReading(
            text=text[:max_chars],
            parts=parts,
            metadata={"format": "docx", "truncated": len(text) > max_chars},
        )
    except OfficeParseError:
        raise
    except Exception as exc:
        raise OfficeParseError(f"DOCX 解析失败：{exc}") from exc


def parse_xlsx(path: Path, max_chars: int = 200_000) -> OfficeReading:
    try:
        from openpyxl import load_workbook  # type: ignore[import-untyped]
    except ImportError as exc:
        raise OfficeParseError("缺少 openpyxl 依赖") from exc
    try:
        workbook = load_workbook(path, read_only=True, data_only=True)
        parts: list[str] = []
        for sheet in workbook.worksheets:
            parts.append(f"# sheet: {sheet.title}")
            for row_index, row in enumerate(sheet.iter_rows(values_only=True)):
                if row_index >= _MAX_ROWS:
                    parts.append("…（超出行数限制，已截断）")
                    break
                values = ["" if value is None else str(value) for value in row[:_MAX_COLS]]
                if any(values):
                    parts.append("\t".join(values))
        text = "\n".join(parts)
        return OfficeReading(
            text=text[:max_chars],
            parts=parts,
            metadata={
                "format": "xlsx",
                "sheets": workbook.sheetnames,
                "truncated": len(text) > max_chars,
            },
        )
    except OfficeParseError:
        raise
    except Exception as exc:
        raise OfficeParseError(f"XLSX 解析失败：{exc}") from exc


def parse_pptx(path: Path, max_chars: int = 200_000) -> OfficeReading:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise OfficeParseError("缺少 python-pptx 依赖") from exc
    try:
        presentation = Presentation(str(path))
        parts: list[str] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            parts.append(f"# slide {slide_index}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
        text = "\n".join(parts)
        return OfficeReading(
            text=text[:max_chars],
            parts=parts,
            metadata={
                "format": "pptx",
                "slides": len(presentation.slides),
                "truncated": len(text) > max_chars,
            },
        )
    except OfficeParseError:
        raise
    except Exception as exc:
        raise OfficeParseError(f"PPTX 解析失败：{exc}") from exc


def parse_office(path: Path, max_chars: int = 200_000) -> OfficeReading:
    suffix = path.suffix.lower()
    if suffix == ".docx":
        return parse_docx(path, max_chars)
    if suffix == ".xlsx":
        return parse_xlsx(path, max_chars)
    if suffix == ".pptx":
        return parse_pptx(path, max_chars)
    raise OfficeParseError(f"不支持的 Office 格式：{suffix}（旧版 .doc/.xls/.ppt 走受控转换器）")
